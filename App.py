#--------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------
#Libraries
import streamlit as st 
import whisper         
import os              
import google.generativeai as GenAI
from pptx import Presentation 
from io import BytesIO
from pptx.util import Pt
from pptx.enum.text import MSO_AUTO_SIZE
import re

#--------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------
# This is the visual part of the page 
st.set_page_config(page_title="Gen", page_icon="🪄")
st.markdown("""
<h1 style="
    color: #FFFFFF;
    text-align: center;
    text-shadow: 2px 2px 10px rgba(0,0,0,0.7);
">
🪄 Transcription and Slide Creator
</h1>
""", unsafe_allow_html=True)

#--------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------
Url_Imagen = "https://i.pinimg.com/originals/cf/a2/39/cfa239195d194b724a9d38362859a1af.jpg"

st.markdown(
    f"""
    <style>
    .stApp {{
        background-image: url("{Url_Imagen}");
        background-size: cover;
        background-position: center;
        background-repeat: no-repeat;
        background-attachment: fixed;
    }}

    /* Capa oscura para que el texto se lea bien */
    .main {{
        background-color: rgba(0, 0, 0, 0.45);
        padding: 20px;
        border-radius: 20px;
    }}
    </style>
    """,
    unsafe_allow_html=True
)

#--------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------
#Enter your API KEY here 
# Asegúrate de que tu secrets.toml tenga la clave correcta
if "API_KEY" in st.secrets:
    GenAI.configure(api_key=st.secrets["API_KEY"])
else:
    st.error("No API Key found in secrets.toml")

#--------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------
# PowerPoint Function
def crear_pptx(texto_generado):
    prs = Presentation()

    # Configuración de colores (puedes cambiarlos aquí)
    COLOR_FONDO = RGBColor(30, 30, 30)    # Gris muy oscuro/negro
    COLOR_TITULO = RGBColor(0, 176, 240)  # Celeste brillante
    COLOR_TEXTO = RGBColor(255, 255, 255) # Blanco

    pattern = r"---\s*SLIDE\s*\d+\s*---\s*(.*?)\s*(?=(?:---\s*SLIDE\s*\d+\s*---)|\Z)"
    slides = re.findall(pattern, texto_generado, flags=re.S)

    if not slides:
        slides = [s for s in re.split(r"---\s*SLIDE", texto_generado) if s.strip()]

    for slide_text in slides:
        lines = [l.strip() for l in slide_text.strip().splitlines() if l.strip()]
        if not lines:
            continue

        title_clean = re.sub(r'^#\s*', '', lines[0]).replace('*', '') 
        
        # Crear diapositiva en blanco para tener control total
        slide_layout = prs.slide_layouts[6] # Layout vacío
        slide = prs.slides.add_slide(slide_layout)

        # 1. Pintar el fondo de la diapositiva
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = COLOR_FONDO

        # 2. Añadir un rectángulo decorativo para el título
        shape = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, Pt(60)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = RGBColor(50, 50, 50) # Gris un poco más claro
        shape.line.fill.background() # Sin borde

        # 3. Añadir el Título
        title_box = slide.shapes.add_textbox(Pt(20), Pt(10), prs.slide_width - Pt(40), Pt(50))
        tf_title = title_box.text_frame
        p_title = tf_title.paragraphs[0]
        p_title.text = title_clean
        p_title.font.bold = True
        p_title.font.size = Pt(32)
        p_title.font.color.rgb = COLOR_TITULO

        # 4. Procesar el contenido (bullets)
        notes_idx = None
        for i, ln in enumerate(lines[1:], start=1):
            if any(ln.lower().startswith(x) for x in ["notes", "notes_slide", "notes:"]):
                notes_idx = i
                break

        bullets_lines = lines[1:notes_idx] if notes_idx else lines[1:]
        
        # Caja de texto para el contenido
        body_box = slide.shapes.add_textbox(Pt(40), Pt(80), prs.slide_width - Pt(80), prs.slide_height - Pt(100))
        tf_body = body_box.text_frame
        tf_body.word_wrap = True

        for i, b in enumerate(bullets_lines):
            # Ignorar la línea de "Visual Concept" en la diapositiva para que no ensucie
            if "visual concept" in b.lower():
                continue
                
            p = tf_body.add_paragraph() if i > 0 else tf_body.paragraphs[0]
            p.text = re.sub(r'^[\*\-\u2022]\s*', '', b)
            p.font.size = Pt(20)
            p.font.color.rgb = COLOR_TEXTO
            p.space_after = Pt(10)

        # 5. Notas del orador
        notes_lines = lines[notes_idx+1:] if notes_idx else bullets_lines
        try:
            slide.notes_slide.notes_text_frame.text = "\n".join(notes_lines)
        except:
            pass

    pptx_io = BytesIO()
    prs.save(pptx_io)
    return pptx_io.getvalue()
#--------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------------
#Transcription Function
Audio_fill = st.file_uploader("Upload your audio so we can transcribe", type=["mp3", "mp4" , "opus" ,"wav", "m4a"])

if Audio_fill is not None:
    st.subheader("🎧Preview your audio")
    st.audio(Audio_fill)

    MAX_FILE_SIZE = 10 * 1024 * 1024
   
    if Audio_fill.size > MAX_FILE_SIZE:
        st.error("The audio is too long or too short. Please upload a file shorter than 3 minutes. (MAX 10MB)")
        # No usamos st.stop() aquí para permitir que la UI siga renderizando si se quiere reiniciar
    else:
       
        with open("temp_audio.mp3", "wb") as f:
            f.write(Audio_fill.getbuffer())
            
        # We show the loading message so the user can wait.
        with st.spinner("Whisper is processing your audio"):
            modelo_whisper = whisper.load_model("base")
            resultado = modelo_whisper.transcribe("temp_audio.mp3")

        st.success("Transcription success")
        with st.expander("Show transcription"):
            st.write(resultado["text"])

        if st.button("✨ Generative Slides"):
            
            with st.spinner("Gemini is creating your slides..."):
              
                modelo_gemini = GenAI.GenerativeModel('models/gemini-2.5-flash') # He actualizado al modelo flash más rápido si está disponible, o usa 1.5
                
                # --- NUEVO PROMPT MEJORADO PARA DISEÑO ---
                instruction = f"""
                Analyze the audio transcript: {resultado['text']} and generate ONLY clearly separated slides following these STRICT rules.

                !!! CRITICAL: LANGUAGE ENFORCEMENT !!!
                1. FIRST, analyze the input text to identify the source language exactly.
                2. YOUR OUTPUT MUST BE 100% IN THAT IDENTIFIED SOURCE LANGUAGE.
                3. IF the audio is in English -> Generate slides/notes in ENGLISH.
                4. IF the audio is in French -> Generate slides/notes in FRENCH.
                5. DO NOT translate to Spanish unless the audio is actually in Spanish.

                === BEGIN DESIGN & CONTENT INSTRUCTIONS ===

                1. GOAL
                Create a visually engaging, well-structured presentation based on the audio.
                Avoid walls of text. Use "Visual Markdown" to make it look professional.

                2. INSTRUCTION DETECTION
                Determine whether the audio contains a clear instruction to create content.

                3. IF A CLEAR INSTRUCTION EXISTS
                Generate a presentation with a MINIMUM of 5 SLIDES.
                Each slide must be clearly separated using the separator below.

                --- SLIDE N ---

                4. SLIDE STRUCTURE (MANDATORY & VISUAL)
                Each slide MUST follow this exact internal structure to ensure it looks organized and colorful:

                # [INSERT RELEVANT EMOJI] TITLE OF THE SLIDE
                
                **Visual Concept:** [Describe in 1 sentence a suggestion for an image or icon that fits this slide, e.g., "A futuristic robot shaking hands with a human"]

                🔹 **[Keyword or Main Idea]:** [Explanation text]
                🔸 **[Keyword or Main Idea]:** [Explanation text]
                🔹 **[Keyword or Main Idea]:** [Explanation text]

                notes_slide:
                Full, natural speaker notes written as if a real presenter were explaining the slide aloud.
                *** THE NOTES MUST BE IN THE SAME LANGUAGE AS THE TRANSCRIPT ***

                5. FORMATTING RULES FOR "PRETTIER" SLIDES
                - Use Emojis (🔹, 🔸, 🚀, 💡, ✅) as bullet points instead of simple dots.
                - ALWAYS bold the key concept at the start of a bullet point (e.g., **Efficiency:**).
                - Keep bullet points concise (maximum 2 lines per point).

                6. IF NO CLEAR INSTRUCTION EXISTS
                Generate ONLY ONE slide stating that an explicit instruction is required (in the source language).
                That slide MUST also include notes_slide.

                7. OUTPUT RESTRICTIONS
                Speaker notes must appear ONLY inside notes_slide.
                Do NOT place notes in the slide body.
                """
                # -----------------------------------------

                answer = modelo_gemini.generate_content(instruction)
                
                st.header("📝 Generated Content")
                
                st.info("Everything is ready! You can review the content below and download your slides.")
                st.write(answer.text)
                
                pptx_data = crear_pptx(answer.text)
                
            
                st.write("") 
                
                st.download_button(
                    label="🚀 DOWNLOAD YOUR POWERPOINT",
                    data=pptx_data,
                    file_name="Presentation.pptx",
                    mime="application/vnd.openxmlformats-officedocument.presentationml.presentation",
                    use_container_width=True 
                )
                st.balloons()

                if os.path.exists("temp_audio.mp3"):
                    os.remove("temp_audio.mp3")


